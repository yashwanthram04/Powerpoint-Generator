import io
import json
import os
import re
import time
import base64
import requests  # <-- Added for AI Pipe support
from typing import Any, Dict, List, Optional

from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.responses import HTMLResponse, StreamingResponse, FileResponse, Response

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------- Optional SDKs ----------
try:
    from openai import OpenAI  # OpenAI-style client (also works for AI Pipe via base_url)
except Exception:
    OpenAI = None

try:
    import anthropic  # Anthropic (Claude)
except Exception:
    anthropic = None

try:
    from google import genai  # Google Gemini
except Exception:
    genai = None

# ---------- App ----------
app = FastAPI(title="Your Text, Your Style – PPTX Generator")

# ---------- Limits / Defaults ----------
MAX_TEXT_CHARS = 60_000
MIN_SLIDES = 10
MAX_SLIDES = 40
MAX_TEMPLATE_BYTES = 30 * 1024 * 1024  # 30 MB
OPENAI_DEFAULT_MODEL = "gpt-4o-mini"        # good default for OpenAI / AI Pipe
ANTHROPIC_DEFAULT_MODEL = "claude-3-5-sonnet-latest"
GEMINI_DEFAULT_MODEL = "gemini-2.5-flash"

# ---------- Front page ----------
@app.get("/", response_class=HTMLResponse)
async def serve_frontend():
    """Serve the main HTML interface"""
    html_path = os.path.join(os.path.dirname(__file__), "index.html")
    try:
        with open(html_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    except FileNotFoundError:
        return HTMLResponse(
            content="<h1>Frontend not found</h1><p>Please ensure index.html is next to app.py</p>",
            status_code=404,
        )

# ---------- Favicon (with fallback) ----------
_FAVICON_FALLBACK_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO3n+9QAAAAASUVORK5CYII="
)

@app.get("/favicon.ico", include_in_schema=False)
async def favicon():
    """
    Serve favicon.ico if present in the working directory.
    Otherwise return a tiny transparent PNG to avoid 404s.
    """
    path = "favicon.ico"
    if os.path.exists(path):
        return FileResponse(path, media_type="image/x-icon")
    return Response(content=_FAVICON_FALLBACK_PNG, media_type="image/png")

# ---------- Generate endpoint ----------
@app.post("/generate")
async def generate_pptx(
    text: str = Form(...),
    guidance: Optional[str] = Form(None),
    provider: str = Form(...),                # "openai", "aipipe", "anthropic", "gemini"
    api_key: str = Form(...),                 # OpenAI key / AI Pipe token / Gemini key / Anthropic key
    model: Optional[str] = Form(None),
    num_slides: Optional[int] = Form(None),   # <-- NEW: desired number of slides (1..40)
    reuse_images: bool = Form(False),         # copy exact images from uploaded PPT slides
    template: Optional[UploadFile] = None,    # OPTIONAL template (.pptx/.potx)
):
    # Validate text
    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Text is required.")

    # Normalize & clamp target slides
    target_slides = None
    if num_slides is not None:
        try:
            target_slides = max(1, min(MAX_SLIDES, int(num_slides)))
        except Exception:
            target_slides = MIN_SLIDES

    # Optional template handling (validate if provided)
    tpl_bytes: Optional[bytes] = None
    if template and template.filename:
        if not template.filename.lower().endswith((".pptx", ".potx")):
            raise HTTPException(status_code=400, detail="Template must be .pptx or .potx.")
        tpl_bytes = await template.read()
        if len(tpl_bytes) < 1024:
            raise HTTPException(status_code=400, detail="Template looks empty or invalid.")
        if len(tpl_bytes) > MAX_TEMPLATE_BYTES:
            raise HTTPException(status_code=400, detail="Template too large (max 30 MB).")

    # Build slide plan using chosen provider (JSON-only)
    try:
        plan = await build_slide_plan_with_retry(
            text=text.strip()[:MAX_TEXT_CHARS],
            guidance=(guidance or "").strip(),
            provider=provider.lower().strip(),
            api_key=api_key.strip(),
            model=(model or "").strip() or None,
            target_slides=target_slides,   # pass through to instruct LLM
            max_retries=2,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM error: {e}")

    # Enforce slide count (exact if provided; otherwise ensure minimum)
    if target_slides:
        plan = enforce_target_slides(plan, target=target_slides, max_slides=MAX_SLIDES)
    else:
        plan = ensure_min_slides(plan, min_slides=MIN_SLIDES, max_slides=MAX_SLIDES)

    # Build PPTX from plan + template style
    try:
        out_bytes = build_presentation_from_plan(
            template_bytes=tpl_bytes,
            plan=plan,
            exact_reuse_images=bool(reuse_images),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PowerPoint build error: {e}")

    headers = {
        "Content-Disposition": 'attachment; filename="generated.pptx"',
        "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "Cache-Control": "no-store",
    }
    return StreamingResponse(io.BytesIO(out_bytes), headers=headers)

# ---------- Slide Plan (LLM) ----------
SLIDE_SCHEMA_EXAMPLE = {
    "slides": [
        {"title": "A short slide title", "bullets": ["point 1", "point 2", "point 3"]}
    ]
}

def _llm_instruction(text: str, guidance: str, target_slides: Optional[int]) -> str:
    count_line = (
        f"- Choose exactly {target_slides} slides (if content is short, expand; if long, summarize).\n"
        if target_slides else
        f"- Choose a reasonable number of slides (min {MIN_SLIDES}, max {MAX_SLIDES})\n"
    )
    return f"""
You are a slide planner. Return JSON ONLY (no code fences, no Markdown), mapping the user's text into slides.

Constraints:
{count_line}- Title ≤ 80 chars
- 3–6 bullets per slide, each ≤ 120 chars
- No images or tables in output
- Omit 'notes' unless essential.

Bias structure & tone toward: "{guidance}"

Return a JSON object like:
{json.dumps(SLIDE_SCHEMA_EXAMPLE, indent=2)}
TEXT:
{text}
""".strip()

async def build_slide_plan_with_retry(
    text: str,
    guidance: str,
    provider: str,
    api_key: str,
    model: Optional[str],
    target_slides: Optional[int],
    max_retries: int = 2,
) -> Dict[str, Any]:
    last_err: Optional[Exception] = None
    for attempt in range(max_retries + 1):
        try:
            return await build_slide_plan(text, guidance, provider, api_key, model, target_slides)
        except Exception as e:
            last_err = e
            if attempt == max_retries:
                break
            time.sleep(0.8 * (attempt + 1))  # small backoff
    assert last_err is not None
    raise last_err

# ---------- AI Pipe call function ----------
def call_aipipe(api_key, model_name, messages, max_tokens=1000):
    """
    Calls AI Pipe's OpenRouter-compatible API for generating completions.
    """
    url = "https://aipipe.org/openrouter/v1/chat/completions"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key.strip()}",
    }

    body = {
        "model": model_name,
        "messages": messages,
        "max_tokens": max_tokens,
        "temperature": 0.7
    }

    response = requests.post(url, headers=headers, json=body)
    if response.status_code != 200:
        raise RuntimeError(f"AI Pipe error {response.status_code}: {response.text}")

    data = response.json()

    # Safely extract message content
    try:
        content = data["choices"][0]["message"]["content"]
        slides = json.loads(content)
        if "slides" not in slides:
            raise ValueError("No slides array in AI Pipe response")
        return slides
    except Exception:
        # Fallback: If AI Pipe sends plain text instead of JSON
        return {"slides": [{"title": "Generated Deck", "bullets": [content]}]}

async def build_slide_plan(
    text: str,
    guidance: str,
    provider: str,
    api_key: str,
    model: Optional[str],
    target_slides: Optional[int],
) -> Dict[str, Any]:
    model_name = (
        model
        or (OPENAI_DEFAULT_MODEL if provider in ("openai", "aipipe") else
            ANTHROPIC_DEFAULT_MODEL if provider == "anthropic" else
            GEMINI_DEFAULT_MODEL)
    )

    instruction = _llm_instruction(text, guidance, target_slides)
    messages = [{"role": "user", "content": instruction}]

    if provider == "aipipe":
        # Use direct AI Pipe call, do NOT use response_format
        return call_aipipe(api_key, model_name, messages)

    if provider == "openai":
        if OpenAI is None:
            raise RuntimeError("openai package not installed. `pip install openai`")
        client_kwargs = {"api_key": api_key}
        client = OpenAI(**client_kwargs)
        resp = client.responses.create(
            model=model_name,
            input=[{"role": "user", "content": instruction}],

            temperature=0.2,
        )
        content = _extract_openai_output_text(resp)
        data = _safe_json_parse(content)

    elif provider == "anthropic":
        if anthropic is None:
            raise RuntimeError("anthropic package not installed. `pip install anthropic`")
        client = anthropic.Anthropic(api_key=api_key)
        msg = client.messages.create(
            model=model_name,
            max_tokens=2048,
            system="Return ONLY valid JSON. No explanations. No code fences.",
            messages=messages,
            temperature=0.2,
        )
        text_out = "".join(
            blk.text for blk in getattr(msg, "content", []) if getattr(blk, "type", "") == "text"
        )
        data = _safe_json_parse(text_out)

    elif provider == "gemini":
        if genai is None:
            raise RuntimeError("google-genai package not installed. `pip install google-genai`")
        client = genai.Client(api_key=api_key)
        resp = client.models.generate_content(
            model=model_name,
            contents=instruction,
            config={"response_mime_type": "application/json"},
        )
        data = _safe_json_parse(resp.text or "")

    else:
        raise HTTPException(status_code=400, detail=f"Unsupported provider: {provider}")

    if not isinstance(data, dict) or "slides" not in data or not isinstance(data["slides"], list):
        raise RuntimeError("Provider did not return JSON with a 'slides' array.")
    return data

def _extract_openai_output_text(resp: Any) -> str:
    txt = getattr(resp, "output_text", None)
    if txt:
        return txt
    try:
        out = getattr(resp, "output", None)
        if out and len(out) and hasattr(out[0], "content") and len(out[0].content):
            maybe = out[0].content[0]
            if hasattr(maybe, "text"):
                return maybe.text
    except Exception:
        pass
    try:
        rs = getattr(resp, "responses", None)
        if rs and len(rs) and hasattr(rs[0], "output_text"):
            return rs[0].output_text
    except Exception:
        pass
    return json.dumps(resp, default=str)

def _safe_json_parse(s: str) -> Dict[str, Any]:
    s = (s or "").strip()
    if not s:
        return {"slides": []}
    try:
        return json.loads(s)
    except Exception:
        m = re.search(r"\{.*\}", s, flags=re.S)
        if m:
            return json.loads(m.group(0))
        raise

# ---------- Enforce slide counts ----------
def ensure_min_slides(plan: Dict[str, Any], min_slides: int, max_slides: int) -> Dict[str, Any]:
    """Ensure the plan has at least `min_slides` (<= max_slides)."""
    slides = plan.get("slides") or []
    out: List[Dict[str, Any]] = []
    for s in slides:
        title = str(s.get("title", "")).strip() or "Slide"
        bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()]
        out.append({"title": title, "bullets": bullets})

    # Split dense slides into chunks of up to 3 bullets
    i = 0
    while len(out) < min_slides and i < len(out) and len(out) < max_slides:
        s = out[i]
        if len(s["bullets"]) > 3:
            extra = s["bullets"][3:]
            s["bullets"] = s["bullets"][:3]
            while extra and len(out) < min_slides and len(out) < max_slides:
                chunk = extra[:3]
                extra = extra[3:]
                out.insert(i + 1, {"title": f"{s['title']} (cont.)", "bullets": chunk})
                i += 1
        i += 1

    # Pad with title-only slides if still fewer than min
    while len(out) < min_slides and len(out) < max_slides:
        out.append({"title": f"Slide {len(out)+1}", "bullets": []})

    plan["slides"] = out[:max_slides]
    return plan

def enforce_target_slides(plan: Dict[str, Any], target: int, max_slides: int) -> Dict[str, Any]:
    """Force the plan to have exactly `target` slides (clamped to max_slides)."""
    target = max(1, min(max_slides, int(target)))
    slides = plan.get("slides") or []
    # Normalize
    norm: List[Dict[str, Any]] = []
    for s in slides:
        norm.append({
            "title": (str(s.get("title", "")) or "Slide").strip(),
            "bullets": [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()],
        })

    # If too few: split & pad
    if len(norm) < target:
        norm_wrap = ensure_min_slides({"slides": norm}, min_slides=target, max_slides=max_slides)["slides"]
        plan["slides"] = norm_wrap[:target]
        return plan

    # If too many: try to merge simple "(cont.)" slides, else truncate
    if len(norm) > target:
        merged: List[Dict[str, Any]] = []
        i = 0
        while i < len(norm):
            cur = norm[i]
            if i + 1 < len(norm) and norm[i + 1]["title"].startswith(cur["title"]):
                # merge a continuation into current (cap bullets ~8)
                nxt = norm[i + 1]
                cur["bullets"].extend(nxt["bullets"])
                cur["bullets"] = cur["bullets"][:8]
                i += 2
                merged.append(cur)
            else:
                merged.append(cur)
                i += 1
        norm = merged

    plan["slides"] = norm[:target]
    return plan

# ---------- PPTX build (no overlap images) ----------
def build_presentation_from_plan(
    template_bytes: Optional[bytes],
    plan: Dict[str, Any],
    exact_reuse_images: bool = False,
) -> bytes:
    """
    Build a new deck:
      - If template provided, inherit its masters/layouts/styles.
      - If not, start from a blank Presentation().
      - If exact_reuse_images=True, copy PICTURE shapes from each template slide
        to the corresponding generated slide, but NEVER cover text:
          • add images BEFORE text (text stays on top)
          • auto-reposition/scale to a safe area if they overlap any text zone
    """
    prs = Presentation(io.BytesIO(template_bytes)) if template_bytes else Presentation()

    # Collect per-slide picture specs BEFORE clearing slides (if requested)
    template_pictures: List[List[Dict[str, any]]] = []
    if template_bytes and exact_reuse_images:
        for s in prs.slides:
            slide_specs: List[Dict[str, any]] = []
            for shape in s.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        slide_specs.append({
                            "blob": shape.image.blob,
                            "left": int(shape.left), "top": int(shape.top),
                            "width": int(shape.width), "height": int(shape.height),
                        })
                    except Exception:
                        pass
            template_pictures.append(slide_specs)

    # Clear existing slides SAFELY (drop relationships first) so Office doesn't “repair”
    _clear_all_slides_safely(prs)

    # Choose a Title+Content layout if possible
    layout_idx = _find_title_and_content_layout_index(prs)
    if layout_idx is None:
        layout_idx = 1 if len(prs.slide_layouts) > 1 else 0

    slides_data = plan.get("slides", [])
    if not isinstance(slides_data, list) or not slides_data:
        raise RuntimeError("No slides found in plan.")

    slide_w, slide_h = int(prs.slide_width), int(prs.slide_height)

    for idx, s in enumerate(slides_data):
        title_txt = str(s.get("title", "")).strip()[:120] or f"Slide {idx+1}"
        bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()][:10]

        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # --- Find placeholder rects early (for safe-zone selection)
        title_rect = None
        body_rect = None
        if slide.shapes.title:
            t = slide.shapes.title
            title_rect = _rect(int(t.left), int(t.top), int(t.width), int(t.height))
        for ph in slide.placeholders:
            try:
                if ph.is_placeholder and ph.placeholder_format.type not in (1,):  # 1 = title
                    body_rect = _rect(int(ph.left), int(ph.top), int(ph.width), int(ph.height))
                    break
            except Exception:
                continue

        # Gather ALL text-bearing zones (titles, subtitles, content placeholders, text boxes, etc.)
        text_zones = _collect_text_zones(slide)

        # --- 1) Insert images FIRST so text is on top (z-order safe) and avoid covering text
        if template_pictures and idx < len(template_pictures):
            for pic in template_pictures[idx]:
                target = _rect(pic["left"], pic["top"], pic["width"], pic["height"])
                if _overlaps_any_text(target, text_zones, thresh=0.10):
                    safe = _choose_safe_zone(
                        slide_w=slide_w, slide_h=slide_h,
                        title_rect=title_rect, body_rect=body_rect,
                        pad=91440  # ~0.1 inch in EMU
                    )
                    target = _fit_into_box(target, safe)
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(pic["blob"]),
                        target["left"], target["top"],
                        width=target["width"], height=target["height"]
                    )
                except Exception:
                    pass

        # --- 2) Now add text so it stays above images
        if slide.shapes.title:
            slide.shapes.title.text = title_txt
        else:
            tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(9), Inches(1))
            tf = tx.text_frame
            tf.text = title_txt
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True

        # Use an existing body/content placeholder if present; else add a textbox
        body_ph = None
        for ph in slide.placeholders:
            try:
                if ph.is_placeholder and ph.placeholder_format.type not in (1,):
                    body_ph = ph
                    break
            except Exception:
                continue

        if body_ph:
            tf = body_ph.text_frame
            tf.clear()
            if bullets:
                tf.paragraphs[0].text = bullets[0]
                tf.paragraphs[0].level = 0
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
            else:
                tf.paragraphs[0].text = ""
        else:
            tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.5), Inches(4.5))
            tf = tx.text_frame
            if bullets:
                tf.paragraphs[0].text = bullets[0]
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
            else:
                tf.paragraphs[0].text = ""

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def _clear_all_slides_safely(prs: Presentation) -> None:
    """Delete all slides and drop relationships to prevent 'repair' prompts in Office."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

def _find_title_and_content_layout_index(prs: Presentation) -> Optional[int]:
    """Heuristic: find a layout that has both a title and a body/content placeholder."""
    for i, layout in enumerate(prs.slide_layouts):
        has_title, has_body = False, False
        try:
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                if t == 1:      # Title
                    has_title = True
                if t in (2, 7): # Body or Content
                    has_body = True
            if has_title and has_body:
                return i
        except Exception:
            continue
    return None

# ---------- Geometry & text-zone helpers ----------
def _collect_text_zones(slide) -> List[Dict[str, int]]:
    """Rects for any shape likely to hold text: title/body/center-title/subtitle/content placeholders,
    plus any autoshape/textbox with a text_frame."""
    zones: List[Dict[str, int]] = []
    for sh in slide.shapes:
        try:
            if getattr(sh, "is_placeholder", False):
                ph_type = getattr(getattr(sh, "placeholder_format", None), "type", None)
                # common text placeholders: 1=TITLE, 2=BODY, 3=CENTER_TITLE, 4=SUBTITLE, 7=CONTENT
                if ph_type in (1, 2, 3, 4, 7):
                    zones.append(_rect(int(sh.left), int(sh.top), int(sh.width), int(sh.height)))
                    continue
            if getattr(sh, "has_text_frame", False):
                zones.append(_rect(int(sh.left), int(sh.top), int(sh.width), int(sh.height)))
        except Exception:
            continue
    return zones

def _rect(left: int, top: int, width: int, height: int) -> Dict[str, int]:
    return {"left": max(0, left), "top": max(0, top), "width": max(0, width), "height": max(0, height)}

def _intersect_area(a: Dict[str, int], b: Dict[str, int]) -> int:
    if not a or not b:
        return 0
    ax1, ay1, ax2, ay2 = a["left"], a["top"], a["left"] + a["width"], a["top"] + a["height"]
    bx1, by1, bx2, by2 = b["left"], b["top"], b["left"] + b["width"], b["top"] + b["height"]
    ix1, iy1, ix2, iy2 = max(ax1, bx1), max(ay1, by1), min(ax2, bx2), min(ay2, by2)
    if ix2 <= ix1 or iy2 <= iy1:
        return 0
    return (ix2 - ix1) * (iy2 - iy1)

def _overlaps_any_text(img: Dict[str, int], zones: List[Dict[str, int]], thresh: float = 0.10) -> bool:
    area = max(1, img["width"] * img["height"])
    for z in zones:
        if _intersect_area(img, z) / area > thresh:
            return True
    return False

def _choose_safe_zone(slide_w: int, slide_h: int,
                      title_rect: Optional[Dict[str, int]],
                      body_rect: Optional[Dict[str, int]],
                      pad: int = 0) -> Dict[str, int]:
    """
    Prefer a column to the RIGHT of the body; if too narrow, use BELOW the body.
    If no body placeholder, fall back to area under the title; else a right sidebar.
    """
    if body_rect:
        # Right of body
        right_left = body_rect["left"] + body_rect["width"] + pad
        right_width = max(0, slide_w - right_left - pad)
        right_top = body_rect["top"]
        right_height = body_rect["height"]
        if right_width >= slide_w * 0.18 and right_height >= slide_h * 0.18:
            return _rect(right_left, right_top, right_width, right_height)

        # Below body
        below_top = body_rect["top"] + body_rect["height"] + pad
        below_height = max(0, slide_h - below_top - pad)
        if below_height >= slide_h * 0.18:
            return _rect(pad, below_top, max(0, slide_w - 2 * pad), below_height)

        # Left of body (last resort)
        left_width = max(0, body_rect["left"] - 2 * pad)
        if left_width >= slide_w * 0.18:
            return _rect(pad, body_rect["top"], left_width, body_rect["height"])

    if title_rect:
        area_top = title_rect["top"] + title_rect["height"] + pad
        area_height = max(0, slide_h - area_top - pad)
        return _rect(pad, area_top, max(0, slide_w - 2 * pad), area_height)

    # Fallback: right sidebar
    sidebar_left = int(slide_w * 0.64) + pad
    sidebar_width = max(0, int(slide_w * 0.36) - 2 * pad)
    sidebar_top = int(slide_h * 0.18) + pad
    sidebar_height = max(0, int(slide_h * 0.72) - 2 * pad)
    return _rect(sidebar_left, sidebar_top, sidebar_width, sidebar_height)

def _fit_into_box(img: Dict[str, int], box: Dict[str, int]) -> Dict[str, int]:
    """Scale img to fit within box, keep aspect ratio; center it inside the box."""
    iw, ih = max(1, img["width"]), max(1, img["height"])
    bw, bh = max(1, box["width"]), max(1, box["height"])
    scale = min(bw / iw, bh / ih, 1.0)
    nw, nh = int(iw * scale), int(ih * scale)
    nl = box["left"] + (bw - nw) // 2
    nt = box["top"] + (bh - nh) // 2
    return _rect(nl, nt, nw, nh)

